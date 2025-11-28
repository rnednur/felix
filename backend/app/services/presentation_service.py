"""
PowerPoint presentation generation from deep research results
Uses Google Gemini 2.0 for content enhancement
"""
import logging
from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import json
import httpx
from datetime import datetime
from app.core.config import settings

logger = logging.getLogger(__name__)


class PresentationService:
    """Generate PowerPoint presentations from research results"""

    # Color schemes
    THEMES = {
        "professional": {
            "primary": RGBColor(31, 78, 121),  # Dark blue
            "secondary": RGBColor(68, 114, 196),  # Light blue
            "accent": RGBColor(237, 125, 49),  # Orange
            "text": RGBColor(51, 51, 51),  # Dark gray
            "background": RGBColor(255, 255, 255),  # White
        },
        "modern": {
            "primary": RGBColor(0, 176, 240),  # Cyan
            "secondary": RGBColor(91, 155, 213),  # Sky blue
            "accent": RGBColor(255, 192, 0),  # Gold
            "text": RGBColor(68, 68, 68),
            "background": RGBColor(255, 255, 255),
        },
        "corporate": {
            "primary": RGBColor(41, 56, 79),  # Navy
            "secondary": RGBColor(79, 129, 189),  # Corporate blue
            "accent": RGBColor(192, 0, 0),  # Red
            "text": RGBColor(64, 64, 64),
            "background": RGBColor(250, 250, 250),  # Off-white
        },
        "vibrant": {
            "primary": RGBColor(142, 68, 173),  # Purple
            "secondary": RGBColor(52, 152, 219),  # Bright blue
            "accent": RGBColor(46, 204, 113),  # Green
            "text": RGBColor(44, 62, 80),
            "background": RGBColor(255, 255, 255),
        }
    }

    def __init__(self):
        self.prs = None
        self.theme = None
        self.gemini_model = settings.PRESENTATION_MODEL

    def _clean_research_data(self, research_result: Dict[str, Any]) -> Dict[str, Any]:
        """
        Remove all Python code, SQL, and technical implementation details
        Keep only meaningful analysis, insights, and findings
        """
        logger.info("🧹 Starting data cleaning (removing code/SQL)...")
        cleaned = research_result.copy()

        # Clean supporting details - remove code/SQL, keep only analysis
        if 'supporting_details' in cleaned:
            supporting_details = cleaned['supporting_details']
            if supporting_details is not None:
                cleaned_details = []
                for detail in supporting_details:
                    if isinstance(detail, dict):
                        cleaned_detail = {
                            'question': detail.get('question', ''),
                            'answer': detail.get('answer', ''),
                            'data': detail.get('data') if not isinstance(detail.get('data'), str) else None,
                        }
                        # Remove code/sql fields
                        cleaned_detail.pop('code', None)
                        cleaned_detail.pop('sql', None)
                        cleaned_detail.pop('python_code', None)
                        cleaned_detail.pop('generated_sql', None)
                        cleaned_detail.pop('generated_code', None)
                        cleaned_details.append(cleaned_detail)
                cleaned['supporting_details'] = cleaned_details
            else:
                cleaned['supporting_details'] = []

        # Clean technical appendix if present
        if 'technical_appendix' in cleaned:
            appendix = cleaned['technical_appendix']
            if isinstance(appendix, dict):
                # Remove all code-related fields
                appendix.pop('queries_executed', None)
                appendix.pop('python_code', None)
                appendix.pop('sql_queries', None)
                appendix.pop('code_snippets', None)
                # Keep only high-level summary
                cleaned['technical_appendix'] = {
                    'methods_summary': appendix.get('methods_summary', 'Data analysis performed using statistical methods')
                }

        # Clean detailed findings
        if 'detailed_findings' in cleaned:
            detailed_findings = cleaned['detailed_findings']
            if detailed_findings is not None:
                cleaned_findings = []
                for finding in detailed_findings:
                    if isinstance(finding, dict):
                        cleaned_finding = {
                            'finding': finding.get('finding', ''),
                            'context': finding.get('context', ''),
                            'impact': finding.get('impact', '')
                        }
                        cleaned_findings.append(cleaned_finding)
                cleaned['detailed_findings'] = cleaned_findings
            else:
                cleaned['detailed_findings'] = []

        logger.info("✅ Data cleaning complete:")
        logger.info(f"   - Supporting details: {len(cleaned.get('supporting_details', []))} items")
        logger.info(f"   - Key findings: {len(cleaned.get('key_findings', []))} items")
        logger.info(f"   - Detailed findings: {len(cleaned.get('detailed_findings', []))} items")
        return cleaned

    async def _enhance_with_gemini(self, cleaned_data: Dict[str, Any]) -> Dict[str, Any]:
        """
        Use Gemini 2.0 to enhance presentation content
        Focuses on creating compelling narratives from the analysis
        """
        logger.info(f"🤖 Starting Gemini enhancement with model: {self.gemini_model}")

        try:
            # Prepare data summary for Gemini
            prompt = f"""You are creating content for a professional business presentation.

Given this research analysis, enhance and structure it for maximum impact in a PowerPoint presentation.

Research Question: {cleaned_data.get('main_question', '')}

Direct Answer: {cleaned_data.get('direct_answer', '')}

Key Findings:
{json.dumps(cleaned_data.get('key_findings', []), indent=2)}

Supporting Details:
{json.dumps(cleaned_data.get('supporting_details', [])[:5], indent=2)}

Your task:
1. Create a compelling executive summary (2-3 sentences)
2. Refine the key findings to be more impactful and business-focused
3. Enhance the direct answer to be more presentation-friendly
4. Suggest 3 actionable recommendations based on the findings

Format your response as JSON:
{{
    "enhanced_executive_summary": "...",
    "enhanced_direct_answer": "...",
    "enhanced_key_findings": ["...", "...", "..."],
    "recommendations": ["...", "...", "..."]
}}

Focus on:
- Clear, concise business language
- Actionable insights
- Compelling narratives
- NO technical jargon or implementation details
"""

            logger.info(f"📝 Prompt length: {len(prompt)} characters")

            request_payload = {
                "model": self.gemini_model,
                "messages": [
                    {"role": "user", "content": prompt}
                ],
                "temperature": 0.7,
                "max_tokens": 2000
            }

            logger.info(f"📤 Sending request to OpenRouter with model: {self.gemini_model}")

            async with httpx.AsyncClient(timeout=60.0) as client:
                response = await client.post(
                    "https://openrouter.ai/api/v1/chat/completions",
                    headers={
                        "Authorization": f"Bearer {settings.OPENROUTER_API_KEY}",
                        "Content-Type": "application/json",
                        "HTTP-Referer": "http://localhost:8000",  # Optional: for OpenRouter analytics
                        "X-Title": "AI Analytics Platform"  # Optional: for OpenRouter analytics
                    },
                    json=request_payload
                )

                logger.info(f"📥 Response status: {response.status_code}")

                if response.status_code != 200:
                    error_detail = response.text
                    logger.error(f"❌ Gemini enhancement failed with status {response.status_code}")
                    logger.error(f"❌ Error response: {error_detail}")
                    logger.warning(f"⚠️ Continuing without Gemini enhancement (will use original content)")
                    return cleaned_data

                result = response.json()
                logger.info(f"✅ Received response from Gemini")

                content = result['choices'][0]['message']['content']
                logger.info(f"📄 Response content length: {len(content)} characters")

                # Parse JSON response
                import re
                json_match = re.search(r'\{.*\}', content, re.DOTALL)
                if json_match:
                    enhanced = json.loads(json_match.group(0))
                    logger.info(f"✅ Successfully parsed JSON response")
                    logger.info(f"📊 Enhancement keys: {list(enhanced.keys())}")

                    # Merge enhancements into cleaned data
                    cleaned_data['executive_summary'] = enhanced.get('enhanced_executive_summary',
                                                                     cleaned_data.get('direct_answer', ''))
                    cleaned_data['direct_answer'] = enhanced.get('enhanced_direct_answer',
                                                                 cleaned_data.get('direct_answer', ''))
                    cleaned_data['key_findings'] = enhanced.get('enhanced_key_findings',
                                                               cleaned_data.get('key_findings', []))
                    cleaned_data['recommendations'] = enhanced.get('recommendations', [])

                    logger.info("✅ Successfully enhanced presentation content with Gemini 2.0")
                    logger.info(f"   - Executive summary: {len(cleaned_data.get('executive_summary', ''))} chars")
                    logger.info(f"   - Key findings: {len(cleaned_data.get('key_findings', []))} items")
                    logger.info(f"   - Recommendations: {len(cleaned_data.get('recommendations', []))} items")
                else:
                    logger.warning("⚠️ Could not parse Gemini response as JSON")
                    logger.warning(f"⚠️ Raw content: {content[:200]}...")

        except httpx.TimeoutException as e:
            logger.error(f"❌ Gemini request timed out: {str(e)}")
            logger.warning("⚠️ Continuing without Gemini enhancement")
        except Exception as e:
            logger.error(f"❌ Gemini enhancement failed: {str(e)}")
            logger.error(f"❌ Exception type: {type(e).__name__}")
            import traceback
            logger.error(f"❌ Traceback: {traceback.format_exc()}")
            logger.warning("⚠️ Continuing without Gemini enhancement")

        return cleaned_data

    async def generate_presentation(
        self,
        research_result: Dict[str, Any],
        theme: str = "professional",
        include_verbose: bool = False,
        output_path: Optional[str] = None
    ) -> str:
        """
        Generate PowerPoint presentation from research results
        Uses Google Gemini 2.0 to enhance content and removes all code/SQL

        Args:
            research_result: Deep research output with findings, etc.
            theme: Color theme (professional, modern, corporate, vibrant)
            include_verbose: Include verbose analysis sections
            output_path: Optional custom output path

        Returns:
            Path to generated PPTX file
        """
        logger.info("=" * 60)
        logger.info(f"🎨 Generating PowerPoint presentation")
        logger.info(f"   Model: {self.gemini_model}")
        logger.info(f"   Theme: {theme}")
        logger.info(f"   Include verbose: {include_verbose}")
        logger.info("=" * 60)

        # Step 1: Clean data - remove all code/SQL
        cleaned_data = self._clean_research_data(research_result)

        # Step 2: Enhance with Gemini 2.0
        enhanced_data = await self._enhance_with_gemini(cleaned_data)
        logger.info("📊 Content preparation complete, starting slide generation...")

        # Initialize presentation
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

        # Set theme
        self.theme = self.THEMES.get(theme, self.THEMES["professional"])

        # Generate slides using enhanced data
        self._add_title_slide(enhanced_data)
        self._add_executive_summary(enhanced_data)
        self._add_key_findings(enhanced_data)
        self._add_supporting_details(enhanced_data)

        if include_verbose and enhanced_data.get("methodology"):
            self._add_methodology(enhanced_data)
            self._add_recommendations(enhanced_data)

        self._add_limitations(enhanced_data)
        self._add_next_steps(enhanced_data)

        # Save presentation
        if not output_path:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = f"data/presentations/research_{timestamp}.pptx"

        import os
        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        self.prs.save(output_path)

        # Get file size
        file_size = os.path.getsize(output_path)
        file_size_kb = file_size / 1024

        logger.info("=" * 60)
        logger.info(f"✅ PowerPoint presentation generated successfully!")
        logger.info(f"   File: {output_path}")
        logger.info(f"   Size: {file_size_kb:.2f} KB")
        logger.info(f"   Slides: {len(self.prs.slides)}")
        logger.info("=" * 60)

        return output_path

    def _add_title_slide(self, data: Dict[str, Any]):
        """Add title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = data.get("main_question", "Research Analysis")
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.theme["primary"]
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(9), Inches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"Deep Research Analysis • {datetime.now().strftime('%B %d, %Y')}"
        subtitle_frame.paragraphs[0].font.size = Pt(20)
        subtitle_frame.paragraphs[0].font.color.rgb = self.theme["secondary"]
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Add decorative line
        line = slide.shapes.add_shape(
            1,  # Line shape
            Inches(2), Inches(5.2), Inches(6), Inches(0)
        )
        line.line.color.rgb = self.theme["accent"]
        line.line.width = Pt(3)

    def _add_executive_summary(self, data: Dict[str, Any]):
        """Add executive summary slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        self._add_slide_title(slide, "Executive Summary")

        # Direct answer
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(2), Inches(8.4), Inches(4.5)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        direct_answer = data.get("direct_answer", "")
        p = text_frame.paragraphs[0]
        p.text = direct_answer
        p.font.size = Pt(18)
        p.font.color.rgb = self.theme["text"]
        p.space_after = Pt(12)
        p.line_spacing = 1.3

    def _add_key_findings(self, data: Dict[str, Any]):
        """Add key findings slides"""
        findings = data.get("key_findings", [])

        # Split into multiple slides if needed (max 5 per slide)
        for i in range(0, len(findings), 5):
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            page_num = (i // 5) + 1
            total_pages = (len(findings) + 4) // 5

            if total_pages > 1:
                title = f"Key Findings ({page_num}/{total_pages})"
            else:
                title = "Key Findings"

            self._add_slide_title(slide, title)

            # Add findings as bullet points
            content_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(2), Inches(8.4), Inches(4.5)
            )
            text_frame = content_box.text_frame
            text_frame.word_wrap = True

            chunk = findings[i:i + 5]
            for idx, finding in enumerate(chunk):
                p = text_frame.add_paragraph() if idx > 0 else text_frame.paragraphs[0]
                p.text = finding if isinstance(finding, str) else finding.get("finding", str(finding))
                p.font.size = Pt(16)
                p.font.color.rgb = self.theme["text"]
                p.space_before = Pt(8)
                p.space_after = Pt(8)
                p.level = 0

                # Add bullet
                p.text = f"• {p.text}"

    def _add_supporting_details(self, data: Dict[str, Any]):
        """Add supporting details slides"""
        details = data.get("supporting_details", [])

        for idx, detail in enumerate(details[:10]):  # Limit to 10 details
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

            # Title
            question = detail.get("question", f"Analysis {idx + 1}")
            self._add_slide_title(slide, question)

            # Answer
            answer_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(2), Inches(8.4), Inches(2)
            )
            text_frame = answer_box.text_frame
            text_frame.word_wrap = True

            answer = detail.get("answer", "")
            if answer:
                p = text_frame.paragraphs[0]
                p.text = answer
                p.font.size = Pt(14)
                p.font.color.rgb = self.theme["text"]
                p.space_after = Pt(10)

            # Add data visualization if available
            if detail.get("data") and isinstance(detail["data"], dict):
                self._add_data_table(slide, detail["data"], top=Inches(4.2))

    def _add_methodology(self, data: Dict[str, Any]):
        """Add methodology slide"""
        methodology = data.get("methodology", {})
        if not methodology:
            return

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_title(slide, "Methodology")

        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(2), Inches(8.4), Inches(4.5)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        # Add methodology sections
        sections = [
            ("Approach", methodology.get("approach")),
            ("Data Sources", methodology.get("data_sources")),
            ("Analysis Methods", methodology.get("analysis_methods"))
        ]

        for idx, (label, content) in enumerate(sections):
            if content:
                p = text_frame.add_paragraph() if idx > 0 else text_frame.paragraphs[0]
                p.text = f"{label}:"
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = self.theme["primary"]
                p.space_before = Pt(12)

                p2 = text_frame.add_paragraph()
                p2.text = content if isinstance(content, str) else str(content)
                p2.font.size = Pt(12)
                p2.font.color.rgb = self.theme["text"]
                p2.space_after = Pt(8)

    def _add_recommendations(self, data: Dict[str, Any]):
        """Add recommendations slide"""
        recommendations = data.get("recommendations", [])
        if not recommendations:
            return

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_title(slide, "Recommendations")

        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(2), Inches(8.4), Inches(4.5)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for idx, rec in enumerate(recommendations):
            p = text_frame.add_paragraph() if idx > 0 else text_frame.paragraphs[0]
            rec_text = rec if isinstance(rec, str) else rec.get("recommendation", str(rec))
            p.text = f"{idx + 1}. {rec_text}"
            p.font.size = Pt(14)
            p.font.color.rgb = self.theme["text"]
            p.space_before = Pt(8)
            p.space_after = Pt(8)

    def _add_limitations(self, data: Dict[str, Any]):
        """Add limitations slide"""
        limitations = data.get("limitations", [])
        if not limitations:
            return

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_title(slide, "Limitations & Caveats")

        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(2), Inches(8.4), Inches(4.5)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for idx, limitation in enumerate(limitations):
            p = text_frame.add_paragraph() if idx > 0 else text_frame.paragraphs[0]
            lim_text = limitation if isinstance(limitation, str) else limitation.get("limitation", str(limitation))
            p.text = f"• {lim_text}"
            p.font.size = Pt(14)
            p.font.color.rgb = self.theme["text"]
            p.space_before = Pt(8)
            p.space_after = Pt(8)

    def _add_next_steps(self, data: Dict[str, Any]):
        """Add next steps slide"""
        follow_ups = data.get("suggested_follow_ups", [])
        if not follow_ups:
            return

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_title(slide, "Suggested Next Steps")

        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(2), Inches(8.4), Inches(4.5)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for idx, follow_up in enumerate(follow_ups):
            p = text_frame.add_paragraph() if idx > 0 else text_frame.paragraphs[0]
            p.text = f"{idx + 1}. {follow_up}"
            p.font.size = Pt(14)
            p.font.color.rgb = self.theme["text"]
            p.space_before = Pt(8)
            p.space_after = Pt(8)

    def _add_slide_title(self, slide, title: str):
        """Add styled title to slide"""
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(1)
        )
        text_frame = title_box.text_frame
        text_frame.text = title
        text_frame.paragraphs[0].font.size = Pt(32)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = self.theme["primary"]

        # Add underline
        line = slide.shapes.add_shape(
            1,  # Line
            Inches(0.5), Inches(1.6), Inches(9), Inches(0)
        )
        line.line.color.rgb = self.theme["accent"]
        line.line.width = Pt(2)

    def _add_data_table(self, slide, data: Dict[str, Any], top: float = Inches(4)):
        """Add simple data table to slide"""
        if not data or not isinstance(data, dict):
            return

        # Convert dict to rows (limited to first 5 items)
        items = list(data.items())[:5]
        if not items:
            return

        rows = len(items) + 1  # +1 for header
        cols = 2

        # Add table
        table = slide.shapes.add_table(
            rows, cols, Inches(0.8), top, Inches(8.4), Inches(2)
        ).table

        # Set column widths
        table.columns[0].width = Inches(4)
        table.columns[1].width = Inches(4.4)

        # Header
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"

        # Style header
        for col in range(cols):
            cell = table.cell(0, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme["primary"]
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)

        # Data rows
        for idx, (key, value) in enumerate(items):
            table.cell(idx + 1, 0).text = str(key)
            table.cell(idx + 1, 1).text = str(value)

            # Style data cells
            for col in range(cols):
                cell = table.cell(idx + 1, col)
                cell.text_frame.paragraphs[0].font.size = Pt(11)
                cell.text_frame.paragraphs[0].font.color.rgb = self.theme["text"]
